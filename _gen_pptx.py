# -*- coding: utf-8 -*-
"""生成業務日誌 LLM 系統開發進度簡報"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy
from pathlib import Path

HYPERLINK_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # 完全空白

# ── 顏色系統 ──────────────────────────────────────────────────
C_DARK   = RGBColor(0x1E, 0x29, 0x3B)   # 深藍黑（背景/標題）
C_PINK   = RGBColor(0xFF, 0x18, 0x5F)   # 品牌粉紅
C_GREEN  = RGBColor(0x22, 0xC5, 0x5E)   # 完成綠
C_YELLOW = RGBColor(0xFF, 0xBF, 0x00)   # 進行中黃
C_GRAY   = RGBColor(0x94, 0xA3, 0xB8)   # 灰（待辦）
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)   # 淡藍背景
C_ACCENT = RGBColor(0x38, 0xBD, 0xF8)   # 天藍重點

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def slide_bg(slide, color=C_DARK):
    bg = add_rect(slide, 0, 0, 13.33, 7.5, fill=color)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

def badge(slide, text, l, t, w=1.2, h=0.35, fill=C_GREEN, tcolor=C_WHITE, size=11):
    add_rect(slide, l, t, w, h, fill=fill)
    add_text(slide, text, l, t, w, h, size=size, bold=True, color=tcolor,
             align=PP_ALIGN.CENTER)

def divider(slide, l, t, w, color=C_PINK, thick=3):
    line = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()

def file_bar(slide, file_links, y=6.82):
    """底部輸出檔案條，支援超連結點擊開檔。
    file_links: list of (display_name, win_path) tuples，
                或純字串（無連結）。
    """
    C_LINK = RGBColor(0x7D, 0xD3, 0xFC)   # sky-300 底線連結
    C_SEP  = RGBColor(0x47, 0x55, 0x69)   # 分隔符灰

    add_rect(slide, 0, y, 13.33, 0.58, fill=RGBColor(0x0A, 0x10, 0x1E))

    tb = slide.shapes.add_textbox(
        Inches(0.25), Inches(y + 0.09), Inches(12.8), Inches(0.42))
    tf = tb.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]

    # 標籤
    r_lbl = para.add_run()
    r_lbl.text = "📁 輸出：  "
    r_lbl.font.size = Pt(10)
    r_lbl.font.bold = True
    r_lbl.font.color.rgb = C_ACCENT

    for i, item in enumerate(file_links):
        # 分隔
        if i > 0:
            r_sep = para.add_run()
            r_sep.text = "   ｜   "
            r_sep.font.size = Pt(10)
            r_sep.font.color.rgb = C_SEP

        # 解包 (name, path) 或純字串
        if isinstance(item, tuple):
            name, fpath = item[0], item[1] if len(item) > 1 else None
        else:
            name, fpath = item, None

        run = para.add_run()
        run.text = name
        run.font.size = Pt(10)

        if fpath:
            url = "file:///" + fpath.replace("\\", "/")
            rPr = run._r.get_or_add_rPr()
            hlinkClick = etree.SubElement(rPr, qn("a:hlinkClick"))
            rel = slide.part.relate_to(url, HYPERLINK_REL, is_external=True)
            hlinkClick.set(qn("r:id"), rel)
            run.font.color.rgb = C_LINK
            run.font.underline = True
        else:
            run.font.color.rgb = C_WHITE

def metric_box(slide, label, value, sub, l, t, w=2.8, fill=RGBColor(0x1E,0x3A,0x52)):
    add_rect(slide, l, t, w, 1.5, fill=fill)
    add_text(slide, label, l+0.15, t+0.1,  w-0.3, 0.35, size=11, color=C_GRAY)
    add_text(slide, value, l+0.15, t+0.38, w-0.3, 0.7,  size=28, bold=True, color=C_WHITE)
    add_text(slide, sub,   l+0.15, t+1.1,  w-0.3, 0.35, size=10, color=C_ACCENT)

# ══════════════════════════════════════════════════════════════
# Slide 1：封面
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 0.12, 7.5, fill=C_PINK)
add_rect(sl, 0.12, 5.8, 13.21, 1.7, fill=RGBColor(0x0F,0x17,0x2A))

add_text(sl, "業務日誌 LLM 智慧萃取系統", 0.6, 1.5, 12, 1.2,
         size=38, bold=True, color=C_WHITE)
add_text(sl, "開發進度 & 成果簡報", 0.6, 2.6, 10, 0.8,
         size=24, color=C_ACCENT)
divider(sl, 0.6, 3.35, 8, color=C_PINK)
add_text(sl, "百萬篇業務日報 → 零人工 L1–L7 分類 → 法人標籤 → 三輸出邏輯",
         0.6, 3.55, 12, 0.6, size=14, color=C_GRAY)
add_text(sl, "2026 年 4 月 23 日　　v4.7", 0.6, 6.0, 8, 0.4,
         size=12, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
# Slide 2：整體架構 & 進度總覽
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "整體進度總覽", 0.5, 0.15, 10, 0.8, size=26, bold=True, color=C_WHITE)
badge(sl, "2026-04-23", 10.8, 0.2, w=2.0, fill=C_PINK)

phases = [
    ("Phase 0",  "GCP 地基 + 問卷 Adapter",              "✅ 完成",  C_GREEN),
    ("Phase L",  "Step 0/1/2/3 全完成 ✅  180萬筆",        "✅ 完成",  C_GREEN),
    ("Phase 1",  "法人標籤 8大類  206,817 家 ✅",          "✅ 完成",  C_GREEN),
    ("Phase M",  "分層聚類 + Map + 熱路徑  206,817 筆 ✅", "✅ 完成",  C_GREEN),
    ("Phase 2",  "日報深度標籤（L1–L7）  執行中 6.8%",     "🔄 執行中", C_YELLOW),
    ("Phase 3",  "痛需累積表 + 三輸出",                   "⬜ 待建立", C_GRAY),
]
for i, (ph, desc, status, sc) in enumerate(phases):
    row = i // 3; col = i % 3
    x = 0.4 + col * 4.3; y = 1.4 + row * 2.6
    add_rect(sl, x, y, 4.0, 2.3, fill=RGBColor(0x1E,0x3A,0x52))
    add_rect(sl, x, y, 4.0, 0.08, fill=sc)
    add_text(sl, ph,   x+0.2, y+0.18, 3.6, 0.5, size=20, bold=True, color=C_WHITE)
    add_text(sl, desc, x+0.2, y+0.65, 3.6, 0.9, size=12, color=C_LIGHT)
    add_rect(sl, x+0.2, y+1.65, 1.6, 0.38, fill=sc)
    add_text(sl, status, x+0.2, y+1.65, 1.6, 0.38, size=11, bold=True,
             color=C_DARK if sc==C_YELLOW else C_WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# Slide 3：框架說明 — L1–L7 分層定義 & 三大信號
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "框架說明　L1–L7 分層定義 & 三大信號", 0.5, 0.15, 12, 0.8,
         size=24, bold=True, color=C_WHITE)

# ── 左側：L1–L7 定義 ──────────────────────────────────────────
add_text(sl, "L 分層框架（業務日報分類維度）", 0.4, 1.18, 7.5, 0.42,
         size=13, bold=True, color=C_ACCENT)
divider(sl, 0.4, 1.58, 7.3, color=C_ACCENT)

layers_def = [
    ("L1", "痛點層",   "客戶現在面臨什麼問題？痛點類型 / 衝擊程度 / 緊迫度"),
    ("L2", "角色層",   "誰在說這件事？決策者 / 壓力來源 / 負責 KPI"),
    ("L3", "目標層",   "客戶想達成什麼？短中長期目標 / 核心 KPI"),
    ("L4", "議題層",   "具體在討論什麼議題？驅動因素 / 客戶立場"),
    ("L5", "評估層",   "正在評估什麼？競爭者 / 採購標準 / 評估項目"),
    ("L6", "方向層",   "決策往哪走？觸發事件 / 策略方向 / 當前溫度"),
    ("L7", "成果層",   "最終結果？成交/流失原因 / 關鍵因素 / 下一步"),
]
LAYER_COLORS = [C_PINK, RGBColor(0xFF,0x6B,0x35), C_YELLOW,
                RGBColor(0x22,0xD3,0xEE), C_ACCENT,
                RGBColor(0xA7,0x8B,0xFA), C_GREEN]
for i, (lyr, name, desc) in enumerate(layers_def):
    y = 1.72 + i * 0.72
    add_rect(sl, 0.4, y, 7.3, 0.62, fill=RGBColor(0x1E,0x3A,0x52))
    add_rect(sl, 0.4, y, 0.08, 0.62, fill=LAYER_COLORS[i])
    add_text(sl, lyr,  0.62, y+0.05, 0.55, 0.52, size=13, bold=True,
             color=LAYER_COLORS[i])
    add_text(sl, name, 1.22, y+0.05, 1.3,  0.28, size=12, bold=True, color=C_WHITE)
    add_text(sl, desc, 1.22, y+0.33, 6.2,  0.25, size=9,  color=C_GRAY)

# ── 右側：三大信號 ─────────────────────────────────────────────
add_text(sl, "問卷三大信號（Phase 0 萃取）", 8.1, 1.18, 5.0, 0.42,
         size=13, bold=True, color=C_ACCENT)
divider(sl, 8.1, 1.58, 4.8, color=C_ACCENT)

signals = [
    (C_PINK,   "痛點信號",   "L1–L4",
     "客戶在問卷中透露的痛點、議題、目標方向",
     "→ 補充 L1/L2/L3/L4 種子語料"),
    (C_YELLOW, "決策角色信號", "L2",
     "填寫問卷的職稱、部門、決策層級",
     "→ 補充 L2 角色種子語料"),
    (C_ACCENT, "決策因素信號", "L5–L7",
     "選廠商標準、在乎的 KPI、評估重點",
     "→ 補充 L5/L6/L7 種子語料"),
]
for i, (color, title, layer_tag, body, impact) in enumerate(signals):
    y = 1.72 + i * 1.78
    add_rect(sl, 8.1, y, 5.0, 1.62, fill=RGBColor(0x1E,0x3A,0x52))
    add_rect(sl, 8.1, y, 5.0, 0.07, fill=color)
    add_rect(sl, 8.1, y+0.12, 0.95, 0.32, fill=color)
    add_text(sl, layer_tag, 8.1,  y+0.12, 0.95, 0.32,
             size=10, bold=True, color=C_DARK if color==C_YELLOW else C_WHITE,
             align=PP_ALIGN.CENTER)
    add_text(sl, title,  9.15, y+0.12, 3.7, 0.35, size=13, bold=True, color=C_WHITE)
    add_text(sl, body,   8.2,  y+0.55, 4.7, 0.5,  size=10, color=C_LIGHT)
    add_text(sl, impact, 8.2,  y+1.2,  4.7, 0.35, size=9,  color=color)

# ── 底部補充說明 ─────────────────────────────────────────────
add_text(sl,
    "問卷 82,105 筆 × 三大信號萃取 → survey_signals.jsonl（407 MB）"
    "   ║   1,802,590 篇日報 × L1–L7 分類 → phase_l_final.csv",
    0.4, 7.1, 12.5, 0.35, size=9, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
# Slide 4：Phase 0 成果
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "Phase 0　GCP 地基 + 問卷 Adapter", 0.5, 0.15, 11, 0.8,
         size=24, bold=True, color=C_WHITE)
badge(sl, "✅ 完成", 11.5, 0.25, fill=C_GREEN)

# ── 目的說明區塊 ───────────────────────────────────────
add_rect(sl, 0.4, 1.12, 12.5, 0.72, fill=RGBColor(0x0F,0x2A,0x3A))
add_rect(sl, 0.4, 1.12, 0.08, 0.72, fill=C_ACCENT)
badge(sl, "目的", 0.55, 1.18, w=0.7, h=0.28, fill=C_ACCENT, tcolor=C_DARK, size=10)
add_text(sl,
    "建立 GCP 資料骨幹（GCS Bucket + UEL Schema）與多源 Adapter，"
    "將百萬篇日報 & 82,105 筆問卷整合為統一事件格式。"
    "問卷三大信號（痛點 / 角色 / 決策因素）萃取後供 Phase L KNN 種子庫補充語料，"
    "並為 Phase 3 痛需累積表提供問卷維度的跨來源比對基礎。",
    1.35, 1.15, 11.3, 0.65, size=11, color=C_LIGHT)

# ── 子任務清單 ─────────────────────────────────────────
items = [
    ("GCS Bucket + UEL Schema",   "✅", "統一事件格式（日報/問卷/客服），含 source_type & source_tags"),
    ("日報 Adapter",               "✅", "法人 ID 對照，time_bucket 注入 → 供 Phase L 使用"),
    ("問卷 Adapter（82,105 筆）",  "✅", "題型分類（問答/多選/單選）+ 三大信號萃取"),
    ("問卷三大信號萃取",            "✅", "痛點(L1–L4) / 決策角色(L2) / 決策因素(L5–L7) → survey_signals.jsonl"),
    ("by-company 匯總",            "✅", "每家公司問卷行為摘要 → survey_by_company.csv（Phase 3 輸入）"),
]
for i, (item, st, note) in enumerate(items):
    y = 1.9 + i * 0.93
    add_rect(sl, 0.4, y, 12.5, 0.82, fill=RGBColor(0x1E,0x3A,0x52))
    add_rect(sl, 0.4, y, 0.08, 0.82, fill=C_GREEN)
    add_text(sl, st,   0.6,  y+0.07, 0.5, 0.68, size=16, bold=True, color=C_GREEN)
    add_text(sl, item, 1.2,  y+0.07, 5.5, 0.35, size=13, bold=True, color=C_WHITE)
    add_text(sl, note, 1.2,  y+0.42, 11.2, 0.3,  size=10, color=C_GRAY)

file_bar(sl, [
    ("survey_signals.jsonl  (407 MB)",    r"D:\yujui\痛點需求地圖\survey_adapter_output\survey_signals.jsonl"),
    ("survey_by_company.csv  (5 MB)",     r"D:\yujui\痛點需求地圖\survey_adapter_output\survey_by_company.csv"),
])

# ══════════════════════════════════════════════════════════════
# Slide 4：Phase L — Step 0 & 1
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "Phase L　Step 0 種子庫 & Step 1 規則快篩", 0.5, 0.15, 11, 0.8,
         size=24, bold=True, color=C_WHITE)
badge(sl, "✅ 完成", 11.5, 0.25, fill=C_GREEN)

# ── 目的說明 ────────────────────────────────────────────
add_rect(sl, 0.4, 1.12, 12.5, 0.72, fill=RGBColor(0x0F,0x2A,0x3A))
add_rect(sl, 0.4, 1.12, 0.08, 0.72, fill=C_ACCENT)
badge(sl, "目的", 0.55, 1.18, w=0.7, h=0.28, fill=C_ACCENT, tcolor=C_DARK, size=10)
add_text(sl,
    "Step 0：從真實日報以 LLM 標注 L1–L7，結合問卷語料，建立高品質向量種子庫（2,375 筆），"
    "作為 Step 2 KNN 分類的比對基準，使後續無需逐筆呼叫 LLM。"
    "  ║  Step 1：關鍵詞規則快篩，直接定案命中 ≥ 2 個關鍵詞的 HIGH 日報（867K 筆），"
    "減少 78% Embedding & KNN 成本。",
    1.35, 1.14, 11.3, 0.65, size=11, color=C_LIGHT)

# ── Step 0 ──────────────────────────────────────────────
add_text(sl, "Step 0　種子庫建立", 0.5, 1.9, 6, 0.5, size=16, bold=True, color=C_ACCENT)
divider(sl, 0.5, 2.35, 5.8, color=C_ACCENT)
s0 = [("LLM 標注",  "27,073 筆有效標注"),
      ("各層種子",  "全層 ≥ 291 筆（L1–L7）"),
      ("KNN 正確率","73.67%（目標 70%）"),
      ("向量模型",  "gemini-embedding-2-preview 3072維")]
for i, (k, v) in enumerate(s0):
    y = 2.48 + i * 0.75
    add_text(sl, k+"：", 0.6, y, 2.5, 0.55, size=12, color=C_GRAY)
    add_text(sl, v,      2.8, y, 4.0, 0.55, size=13, bold=True, color=C_WHITE)

# ── Step 1 ──────────────────────────────────────────────
add_text(sl, "Step 1　規則快篩", 7.0, 1.9, 6, 0.5, size=16, bold=True, color=C_ACCENT)
divider(sl, 7.0, 2.35, 5.8, color=C_ACCENT)
metric_box(sl, "總處理筆數",  "4,141,953", "篇業務日報", 7.0, 2.48, w=2.7)
metric_box(sl, "HIGH 直接定案", "867,023", "20.9%", 9.9, 2.48, w=2.7)
s1_detail = [("MEDIUM（待 KNN）", "961,849 筆  23.2%"),
             ("SKIP（短文本/雜訊）", "2,313,081 筆  55.8%"),
             ("HIGH 門檻",          "命中關鍵詞 ≥ 2 個")]
for i, (k, v) in enumerate(s1_detail):
    y = 4.6 + i * 0.68
    add_text(sl, k+"：", 7.0, y, 3.2, 0.55, size=12, color=C_GRAY)
    add_text(sl, v,      9.8, y, 3.5, 0.55, size=12, bold=True, color=C_WHITE)

file_bar(sl, [
    ("labeled_logs.jsonl  (27,073筆)",     r"D:\yujui\痛點需求地圖\seed_output\labeled_logs.jsonl"),
    ("seed_vectors.jsonl  (2,375筆)",      r"D:\yujui\痛點需求地圖\seed_output\seed_vectors.jsonl"),
    ("quality_report.csv",                 r"D:\yujui\痛點需求地圖\seed_output\quality_report.csv"),
    ("step1_results.jsonl  (867,023筆)",   r"D:\yujui\痛點需求地圖\step1_output\step1_results.jsonl"),
])

# ══════════════════════════════════════════════════════════════
# Slide 5：Phase L — Step 2 成果
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "Phase L　Step 2　Embedding + KNN 分類", 0.5, 0.15, 11, 0.8,
         size=24, bold=True, color=C_WHITE)
badge(sl, "✅ 完成", 11.5, 0.25, fill=C_GREEN)

# ── 目的說明 ────────────────────────────────────────────
add_rect(sl, 0.4, 1.12, 12.5, 0.65, fill=RGBColor(0x0F,0x2A,0x3A))
add_rect(sl, 0.4, 1.12, 0.08, 0.65, fill=C_ACCENT)
badge(sl, "目的", 0.55, 1.18, w=0.7, h=0.28, fill=C_ACCENT, tcolor=C_DARK, size=10)
add_text(sl,
    "將 Step 1 未定案的 MEDIUM 日報（961,849 筆）以 Gemini Embedding 向量化後做 KNN 比對，"
    "相似度 ≥ 0.65 → HIGH 直接定案，0.50–0.65 → 送 Step 3 精分，< 0.50 → 放棄。"
    "目標：以零 LLM 費用完成 95%+ 的日報 L 分類，將剩餘模糊案例壓到最小。",
    1.35, 1.14, 11.3, 0.6, size=11, color=C_LIGHT)

# ── 數字指標 ─────────────────────────────────────────────
metric_box(sl, "總分類筆數", "935,570", "篇日報", 0.4, 1.85, w=3.0)
metric_box(sl, "HIGH（≥0.65）", "935,556", "100.0%", 3.6, 1.85, w=3.0)
metric_box(sl, "MEDIUM", "14", "0.0%　→ Step 3", 6.8, 1.85, w=3.0)
metric_box(sl, "uncertain", "36", "seed 無對應", 10.0, 1.85, w=3.0)

layers = [("L1","343,184","36.7%"), ("L2","176,434","18.9%"),
          ("L3","114,686","12.3%"), ("L4", "25,695", "2.7%"),
          ("L5","119,093","12.7%"), ("L6","107,900","11.5%"), ("L7","48,528","5.2%")]
bar_total = 935556
BAR_BASE_S2 = 5.55   # 往下移 0.65
for i, (l, n, p) in enumerate(layers):
    x = 0.4 + i * 1.84
    bar_h = int(n.replace(",","")) / bar_total * 2.0
    by = BAR_BASE_S2 - bar_h
    add_rect(sl, x, by, 1.5, bar_h, fill=C_PINK)
    add_text(sl, l,  x, BAR_BASE_S2+0.05, 1.5, 0.32, size=12, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, n,  x, by-0.38, 1.5, 0.35, size=10,
             color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(sl, p,  x, by-0.68, 1.5, 0.3,  size=9,
             color=C_GRAY, align=PP_ALIGN.CENTER)

add_text(sl, "DEBUG：JSONL 無換行（raw_decode 修復）、step2_result 全 uncertain（SEED_LAYER lookup 修補 935,534 筆）",
         0.4, 6.18, 12.5, 0.5, size=9, color=RGBColor(0x64,0x74,0x8B))

file_bar(sl, [
    ("step2_results.jsonl  (935,570筆)",  r"D:\yujui\痛點需求地圖\step2_output\step2_results.jsonl"),
    ("_a2_matched.jsonl  (334.3 MB)",     r"D:\yujui\痛點需求地圖\step2_output\_a2_matched.jsonl"),
    ("step2_vectors.parquet",             r"D:\yujui\痛點需求地圖\step2_output\step2_vectors.parquet"),
])

# ══════════════════════════════════════════════════════════════
# Slide 6：Phase L — Step 3（待執行）
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "Phase L　Step 3　Self-Consistency LLM", 0.5, 0.15, 11, 0.8,
         size=24, bold=True, color=C_WHITE)
badge(sl, "✅ 完成", 11.5, 0.25, fill=C_GREEN)

# ── 目的說明 ────────────────────────────────────────────
add_rect(sl, 0.4, 1.12, 12.5, 0.62, fill=RGBColor(0x0F,0x2A,0x3A))
add_rect(sl, 0.4, 1.12, 0.08, 0.62, fill=C_ACCENT)
badge(sl, "目的", 0.55, 1.18, w=0.7, h=0.28, fill=C_ACCENT, tcolor=C_DARK, size=10)
add_text(sl,
    "對 Step 2 仍模糊的 MEDIUM 日報（14 筆），同一篇送 Gemini 三次投票取多數決，"
    "消除單次 LLM 誤差；3 票一致 → CONFIRM_HIGH，2 票一致 → CONFIRM_MED，不一致 → UNCERTAIN 降權。"
    "三步驟結果合併為 phase_l_final.csv，完成 180 萬筆零人工 L 分類。",
    1.35, 1.14, 11.3, 0.57, size=11, color=C_LIGHT)

# ── 數字指標 ─────────────────────────────────────────────
metric_box(sl, "輸入（MEDIUM）",     "14",   "筆", 0.4, 1.82, w=3.0)
metric_box(sl, "CONFIRM_HIGH(3/3)", "10",   "71.4%", 3.6, 1.82, w=3.0)
metric_box(sl, "CONFIRM_MED(2/3)",  "1",    "7.1%", 6.8, 1.82, w=3.0)
metric_box(sl, "UNCERTAIN",         "3",    "21.4%  降權", 10.0, 1.82, w=3.0)

add_rect(sl, 0.4, 3.4, 12.5, 1.4, fill=RGBColor(0x1E,0x3A,0x52))
add_text(sl, "最終合併　phase_l_final.csv", 0.7, 3.5, 8, 0.5,
         size=14, bold=True, color=C_ACCENT)
add_text(sl, "Step1 HIGH 867,023  +  Step2 HIGH 935,556  +  Step3 CONFIRM 11  =  1,802,590 筆",
         0.7, 3.95, 12.0, 0.45, size=13, color=C_WHITE)

layers_f = [("L1","461,267"), ("L2","366,442"), ("L3","183,624"), ("L4","56,679"),
            ("L5","433,565"), ("L6","114,567"), ("L7","186,410")]
bar_total_f = 1802590
BAR_BASE = 6.05   # 長條圖底線
MAX_BAR  = 1.3    # 最大高度
for i, (l, n) in enumerate(layers_f):
    x = 0.4 + i * 1.84
    bar_h = int(n.replace(",","")) / bar_total_f * MAX_BAR
    by = BAR_BASE - bar_h
    add_rect(sl, x, by, 1.5, bar_h, fill=C_GREEN)
    add_text(sl, l,  x, BAR_BASE+0.05, 1.5, 0.3, size=11, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, n,  x, by-0.35, 1.5, 0.32, size=9,
             color=C_ACCENT, align=PP_ALIGN.CENTER)

file_bar(sl, [
    ("step3_results.jsonl  (14筆)",          r"D:\yujui\痛點需求地圖\step3_output\step3_results.jsonl"),
    ("phase_l_final.csv  (1,802,590筆) ★",  r"D:\yujui\痛點需求地圖\step3_output\phase_l_final.csv"),
], y=6.6)

# ══════════════════════════════════════════════════════════════
# Slide 7：Phase 1 法人標籤
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "Phase 1　法人歷程標籤（8 大類）", 0.5, 0.15, 11, 0.8,
         size=24, bold=True, color=C_WHITE)
badge(sl, "✅ 完成", 11.5, 0.25, fill=C_GREEN)

metric_box(sl, "已完成", "206,817", "家", 0.4, 1.2, w=4.0)
metric_box(sl, "有萃取資料", "106,868", "51.7%", 4.6, 1.2, w=4.0)
metric_box(sl, "空結果（正常）", "99,949", "日報資訊不足", 8.8, 1.2, w=4.0)

cats = ["基本資料", "聯絡人", "進出口資訊", "銷售資訊", "集團資訊", "偏好及總結", "產業", "產業鏈"]
fields = [
    "外商 / 家族企業 / 商業模式 / 員工人數",
    "決策者 / 家族企業人員",
    "國家",
    "競爭者",
    "集團 / 母公司 / 關係企業 / 分公司 / 品牌名稱",
    "關注議題（限5個）",
    "產業龍頭",
    "法人的客戶 / 法人的供應商",
]
for i, (cat, fld) in enumerate(zip(cats, fields)):
    row = i // 4; col = i % 4
    x = 0.4 + col * 3.2; y = 2.95 + row * 1.6
    add_rect(sl, x, y, 3.0, 1.4, fill=RGBColor(0x1E,0x3A,0x52))
    add_rect(sl, x, y, 3.0, 0.07, fill=C_PINK)
    add_text(sl, cat, x+0.15, y+0.12, 2.7, 0.45, size=13, bold=True, color=C_WHITE)
    add_text(sl, fld, x+0.15, y+0.55, 2.7, 0.75, size=10, color=C_GRAY)

add_text(sl, "model: gemini-2.0-flash　RPM=30　6 Workers　RESUME=True　全部完成",
         0.4, 6.3, 12.5, 0.42, size=10, color=C_GRAY)

file_bar(sl, [
    ("company_labels.jsonl  (206,817筆)",  r"D:\yujui\痛點需求地圖\corp_label_output\company_labels.jsonl"),
    ("company_labels_flat.csv",            r"D:\yujui\痛點需求地圖\corp_label_output\company_labels_flat.csv"),
])

# ══════════════════════════════════════════════════════════════
# Slide 8：Phase M 分層聚類成果
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "Phase M　分層聚類 + Map + 熱路徑", 0.5, 0.15, 11, 0.8,
         size=24, bold=True, color=C_WHITE)
badge(sl, "✅ 完成", 11.5, 0.25, fill=C_GREEN)

# 目的
add_rect(sl, 0.4, 1.12, 12.5, 0.62, fill=RGBColor(0x0F,0x2A,0x3A))
add_rect(sl, 0.4, 1.12, 0.08, 0.62, fill=C_ACCENT)
badge(sl, "目的", 0.55, 1.18, w=0.7, h=0.28, fill=C_ACCENT, tcolor=C_DARK, size=10)
add_text(sl,
    "以法人標籤 15 個特徵（外商/商業模式/員工數/貿易國數…）進行 KMeans(k=7) 分群，"
    "PCA 降維視覺化後分析各群 L 層痛需熱路徑，找出不同類型法人的核心需求模式。",
    1.35, 1.14, 11.3, 0.57, size=11, color=C_LIGHT)

# 指標（小型）
for mi, (lbl, val, sub, mx) in enumerate([
    ("分群對象", "206,817 家", "法人", 0.4),
    ("分群數",   "7 clusters", "KMeans", 3.6),
    ("特徵維度", "15 維",      "法人標籤特徵", 6.8),
    ("熱路徑",   "Top 30",    "L 層序列", 10.0),
]):
    add_rect(sl, mx, 1.82, 3.0, 0.85, fill=RGBColor(0x1E,0x3A,0x52))
    add_text(sl, lbl, mx+0.15, 1.87, 2.7, 0.28, size=10, color=C_GRAY)
    add_text(sl, val, mx+0.15, 2.14, 2.7, 0.42, size=16, bold=True, color=C_WHITE)

# 流程：(step, 標題, 目的, 結果)
steps_m = [
    ("A1",    "特徵萃取",
     "目的：將法人標籤 JSON 轉為可計算的數值向量（布林/等級/計數）",
     "結果：15 維特徵 — 外商/家族/商業模式/員工數/決策者數/貿易國數/競爭者數/集團/子公司/品牌/議題/產業龍頭/客戶數/供應商數"),
    ("A2",    "L 層路徑索引",
     "目的：建立每家公司跨月份的 L 層出現記錄，作為熱路徑分析的資料基礎",
     "結果：company_index = {company_id → {L層 → [年月列表]}}，供 C1/C2 步驟使用"),
    ("B1–B3", "標準化 → 分群 → 視覺化",
     "目的：B1 消除量綱差異；B2 KMeans(k=7) 將 206,817 家依特徵分成 7 個客群；B3 PCA 降至 2 維可視覺化",
     "結果：cluster 0–6 各代表一種法人類型（如：外商製造 / 中小內銷 / 集團型…），pca_x/y 供散佈圖呈現"),
    ("C1–C2", "壓縮路徑 → 熱路徑統計",
     "目的：C1 compress_path() 去除連續重複 L 層（L1→L1→L3 → L1→L3），萃取關鍵需求轉折；C2 統計各群 Top-30 路徑",
     "結果：hotpaths.csv — 最常見路徑（如 L1→L3 佔 23%、L2→L5→L7 佔 18%），顯示各類法人的需求演進模式"),
    ("E2",    "輸出三份結果檔",
     "目的：將分群結果、群體特徵摘要、熱路徑彙整為可分析的 CSV",
     "結果：company_clusters.csv（每家公司+cluster+pca座標+15特徵）/ cluster_profiles.csv（7群均值）/ hotpaths.csv（Top-30路徑）"),
]
for i, (step, title, purpose, result) in enumerate(steps_m):
    y = 2.82 + i * 0.88
    add_rect(sl, 0.4, y, 12.5, 0.80, fill=RGBColor(0x1E,0x3A,0x52))
    add_rect(sl, 0.4, y, 0.08, 0.80, fill=C_GREEN)
    add_text(sl, step,    0.60, y+0.04, 1.1,  0.30, size=11, bold=True, color=C_ACCENT)
    add_text(sl, title,   1.75, y+0.04, 3.5,  0.30, size=11, bold=True, color=C_WHITE)
    add_text(sl, purpose, 1.75, y+0.34, 10.9, 0.22, size=8.5, color=C_GRAY)
    add_text(sl, result,  1.75, y+0.55, 10.9, 0.22, size=8.5, color=C_ACCENT)

file_bar(sl, [
    ("company_clusters.csv  (206,817筆)",  r"D:\yujui\痛點需求地圖\phaseM_output\company_clusters.csv"),
    ("cluster_profiles.csv",              r"D:\yujui\痛點需求地圖\phaseM_output\cluster_profiles.csv"),
    ("hotpaths.csv",                      r"D:\yujui\痛點需求地圖\phaseM_output\hotpaths.csv"),
])

_PCA_IMG   = Path(r"D:\yujui\痛點需求地圖\phaseM_output\cluster_pca_map.png")
_RADAR_IMG = Path(r"D:\yujui\痛點需求地圖\phaseM_output\cluster_radar.png")

# ══════════════════════════════════════════════════════════════
# Slide 9：Phase M — PCA 分群地圖
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.0, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "Phase M　PCA 分群地圖", 0.5, 0.13, 10, 0.75,
         size=24, bold=True, color=C_WHITE)
badge(sl, "✅ 完成", 11.5, 0.2, fill=C_GREEN)
add_text(sl, "每個點代表一家法人（抽樣 40,000 家），顏色代表所屬 cluster；標籤為群心位置",
         0.4, 1.02, 12.5, 0.32, size=10, color=C_GRAY)
if _PCA_IMG.exists():
    sl.shapes.add_picture(str(_PCA_IMG), Inches(0.4), Inches(1.35),
                          width=Inches(12.5), height=Inches(5.8))
file_bar(sl, [
    ("cluster_pca_map.png",       r"D:\yujui\痛點需求地圖\phaseM_output\cluster_pca_map.png"),
    ("company_clusters.csv",      r"D:\yujui\痛點需求地圖\phaseM_output\company_clusters.csv"),
])

# ══════════════════════════════════════════════════════════════
# Slide 10：Phase M — 各群特徵雷達圖
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.0, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "Phase M　各 Cluster 特徵雷達圖", 0.5, 0.13, 10, 0.75,
         size=24, bold=True, color=C_WHITE)
badge(sl, "✅ 完成", 11.5, 0.2, fill=C_GREEN)
add_text(sl, "雷達圖面積越大 → 該維度越突出；8 個特徵維度：外商 / 商業模式 / 員工規模 / 貿易國 / 競爭者 / 品牌數 / 議題數 / 客戶數",
         0.4, 1.02, 12.5, 0.32, size=10, color=C_GRAY)
if _RADAR_IMG.exists():
    sl.shapes.add_picture(str(_RADAR_IMG), Inches(0.4), Inches(1.35),
                          width=Inches(12.5), height=Inches(5.8))
file_bar(sl, [
    ("cluster_radar.png",         r"D:\yujui\痛點需求地圖\phaseM_output\cluster_radar.png"),
    ("cluster_profiles.csv",      r"D:\yujui\痛點需求地圖\phaseM_output\cluster_profiles.csv"),
])

# ══════════════════════════════════════════════════════════════
# Slide 10：Phase 2 日報深度標籤（執行中）
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "Phase 2　日報深度標籤（L1–L7 結構化）", 0.5, 0.15, 11, 0.8,
         size=24, bold=True, color=C_WHITE)
badge(sl, "🔄 執行中", 11.0, 0.25, w=2.0, fill=C_YELLOW, tcolor=C_DARK)

# 目的
add_rect(sl, 0.4, 1.12, 12.5, 0.62, fill=RGBColor(0x0F,0x2A,0x3A))
add_rect(sl, 0.4, 1.12, 0.08, 0.62, fill=C_ACCENT)
badge(sl, "目的", 0.55, 1.18, w=0.7, h=0.28, fill=C_ACCENT, tcolor=C_DARK, size=10)
add_text(sl,
    "對每家公司所有 L 層業務日報，一次 LLM 呼叫萃取各層結構化深度標籤："
    "L1 痛點類型/衝擊/緊迫度、L2 角色/壓力來源/KPI、L3 目標/時程/核心KPI…等，"
    "供 Phase 3 痛需累積表與三輸出邏輯使用。",
    1.35, 1.14, 11.3, 0.57, size=11, color=C_LIGHT)

# 進度指標
metric_box(sl, "已完成", "12,165", "家（RESUME）", 0.4, 1.82, w=3.0)
metric_box(sl, "總目標", "178,790", "家", 3.6, 1.82, w=3.0)
metric_box(sl, "進度", "6.8%", "持續執行中", 6.8, 1.82, w=3.0)
metric_box(sl, "RPM 限制", "30", "6 Workers", 10.0, 1.82, w=3.0)

# 各 L 層萃取欄位（中文 + 範例）
# (層, 層名, 萃取欄位中文, 範例值)
layer_fields = [
    ("L1", "痛點層",
     "痛點類型 / 衝擊程度 / 緊迫度",
     "例：原料成本上漲 / 高 / 急迫"),
    ("L2", "角色層",
     "決策角色 / 壓力來源 / 負責KPI",
     "例：採購經理 / 總部壓成本 / 毛利率"),
    ("L3", "目標層",
     "目標類型 / 時程 / 核心KPI",
     "例：降低庫存 / Q3達成 / 周轉天數"),
    ("L4", "議題層",
     "議題名稱 / 驅動因素 / 客戶立場",
     "例：ERP升級 / 法規要求 / 觀望評估"),
    ("L5", "評估層",
     "評估項目 / 競爭者 / 決策標準",
     "例：交期/品質 / SAP/Oracle / 價格優先"),
    ("L6", "方向層",
     "策略方向 / 觸發事件 / 當前溫度",
     "例：轉向數位化 / 客訴增加 / 積極"),
    ("L7", "成果層",
     "結果類型 / 關鍵因素 / 下一步",
     "例：成交 / 價格合理 / 安排簽約"),
]
# 7 格排成 4+3，每格寬 3.1
for i, (layer, lname, fields, example) in enumerate(layer_fields):
    col = i % 4; row2 = i // 4
    x = 0.25 + col * 3.22; y = 3.45 + row2 * 1.62
    add_rect(sl, x, y, 3.05, 1.50, fill=RGBColor(0x1E,0x3A,0x52))
    add_rect(sl, x, y, 3.05, 0.07, fill=C_YELLOW)
    # 層標籤 + 層名
    add_text(sl, layer,  x+0.12, y+0.10, 0.55, 0.38, size=13, bold=True, color=C_YELLOW)
    add_text(sl, lname,  x+0.72, y+0.10, 2.2,  0.38, size=13, bold=True, color=C_WHITE)
    # 欄位
    add_text(sl, fields,  x+0.12, y+0.50, 2.85, 0.32, size=9.5, color=C_LIGHT)
    # 範例
    add_text(sl, example, x+0.12, y+0.88, 2.85, 0.52, size=8.5, color=C_GRAY)

add_text(sl, "model: gemini-2.0-flash　RPM=30　6 Workers　RESUME=True",
         0.4, 6.3, 12.5, 0.42, size=10, color=C_GRAY)

file_bar(sl, [
    ("phase2_deep_labels.jsonl  (執行中)",  r"D:\yujui\痛點需求地圖\phase2_output\phase2_deep_labels.jsonl"),
    ("phase2_labels_flat.csv",             r"D:\yujui\痛點需求地圖\phase2_output\phase2_labels_flat.csv"),
])

# ══════════════════════════════════════════════════════════════
# Slide 10：下一步行動
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=RGBColor(0x0F,0x17,0x2A))
add_text(sl, "下一步行動", 0.5, 0.15, 10, 0.8, size=26, bold=True, color=C_WHITE)

C_LINK = RGBColor(0x7D, 0xD3, 0xFC)

nexts = [
    (C_YELLOW, "執行中",  "Phase 2 日報深度標籤",
               "RESUME=True 持續執行中（6.8%）。完成後跑 E1→E2 輸出 CSV",
               [("phase2_deep_labels.jsonl  (12,165筆)",  r"D:\yujui\痛點需求地圖\phase2_output\phase2_deep_labels.jsonl"),
                ("phase2_labels_flat.csv",               r"D:\yujui\痛點需求地圖\phase2_output\phase2_labels_flat.csv")]),
    (C_PINK,   "完成後",  "Phase 3 痛需累積表",
               "彙整 Phase L + Phase 1 + Phase 2 → 公司 × 痛需維度累積表，三輸出邏輯（推方案/問項/開發卡）",
               [("phase3_output/ (待建立)", None)]),
    (C_ACCENT, "後續",   "Phase 3 三輸出邏輯",
               "推薦方案、標準問項庫、開發機會卡 → 銷售 AI 助理輸入",
               [("phase3_output/ (待建立)", None)]),
    (C_GRAY,   "未來",   "Phase 4+ 整合應用",
               "結合所有輸出，建立業務 AI copilot（推薦 / 提問 / 預警）",
               [("TBD", None)]),
]

def add_file_links_inline(slide, file_links, l, t, w, h):
    """在文字框中插入可點擊的超連結檔案列表"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    r_pre = para.add_run()
    r_pre.text = "→ "
    r_pre.font.size = Pt(9)
    r_pre.font.color.rgb = C_GRAY
    for i, item in enumerate(file_links):
        if i > 0:
            r_sep = para.add_run()
            r_sep.text = "  /  "
            r_sep.font.size = Pt(9)
            r_sep.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
        name = item[0] if isinstance(item, tuple) else item
        fpath = item[1] if isinstance(item, tuple) else None
        run = para.add_run()
        run.text = name
        run.font.size = Pt(9)
        if fpath:
            url = "file:///" + fpath.replace("\\", "/")
            rPr = run._r.get_or_add_rPr()
            hlinkClick = etree.SubElement(rPr, qn("a:hlinkClick"))
            rel = slide.part.relate_to(url, HYPERLINK_REL, is_external=True)
            hlinkClick.set(qn("r:id"), rel)
            run.font.color.rgb = C_LINK
            run.font.underline = True
        else:
            run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

for i, (color, timing, title, desc, file_links) in enumerate(nexts):
    y = 1.2 + i * 1.45
    add_rect(sl, 0.4, y, 12.5, 1.3, fill=RGBColor(0x1E,0x3A,0x52))
    add_rect(sl, 0.4, y, 0.1, 1.3, fill=color)
    add_rect(sl, 0.65, y+0.13, 1.2, 0.38, fill=color)
    add_text(sl, timing, 0.65, y+0.13, 1.2, 0.38, size=10, bold=True,
             color=C_DARK if color==C_YELLOW else C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 2.0, y+0.1,  5.5, 0.45, size=14, bold=True, color=C_WHITE)
    add_text(sl, desc,  2.0, y+0.55, 10.5, 0.35, size=10, color=C_GRAY)
    add_file_links_inline(sl, file_links, 2.0, y+0.9, 10.5, 0.32)

DST = r"d:\yujui\痛點需求地圖\prompt定版\LLM系統開發進度簡報_v5_0.pptx"
prs.save(DST)
print("OK:", DST)
print(f"共 {len(prs.slides)} 張投影片")
